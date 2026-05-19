# make_data.py
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

os.makedirs("data/raw", exist_ok=True)

# 1. 產生 PDF 規格書
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", 'B', 16)
pdf.cell(0, 10, "Raydium AI Display Driver Spec", 0, 1)
pdf.ln(10)
pdf.set_font("Arial", size=12)
pdf.multi_cell(0, 10, "Model: RD-AI-2026\nOperating Voltage: 0.8V - 1.2V\n"
                      "Interface: MIPI DSI-2\nMax Resolution: 3200x1440 @ 90fps\n"
                      "Features: Embedded hardware super-resolution pipeline.")
pdf.output("data/raw/specs.pdf")

# 2. 產生 PPTX 簡報
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
box.text_frame.text = "Raydium Performance Metrics: Power consumption is 45 mW at FHD+."
prs.save("data/raw/presentation.pptx")

print("✅ 階段 1 完成：模擬技術文件已生成於 data/raw/")