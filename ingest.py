# ingest.py
import os
import json
from pypdf import PdfReader
from pptx import Presentation

def chunk_text(text, chunk_size=150, chunk_overlap=30):
    chunks = []
    text = " ".join(text.split())
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += (chunk_size - chunk_overlap)
    return chunks

def process_files():
    raw_dir = "data/raw"
    processed_data = []
    
    if not os.path.exists(raw_dir): return

    for filename in os.listdir(raw_dir):
        file_path = os.path.join(raw_dir, filename)
        content = ""
        if filename.endswith(".pdf"):
            reader = PdfReader(file_path)
            content = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif filename.endswith(".pptx"):
            prs = Presentation(file_path)
            content = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        
        if content:
            chunks = chunk_text(content)
            for idx, chunk in enumerate(chunks):
                processed_data.append({
                    "source": filename,
                    "chunk_id": f"{filename}_chunk_{idx}",
                    "content": chunk
                })

    os.makedirs("data/processed", exist_ok=True)
    with open("data/processed/knowledge_base.json", "w", encoding="utf-8") as f:
        json.dump(processed_data, f, ensure_ascii=False, indent=4)
    print(f"✅ 階段 2 完成：成功切出 {len(processed_data)} 個知識區塊！")

if __name__ == "__main__":
    process_files()